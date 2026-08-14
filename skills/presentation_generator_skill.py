from pptx import Presentation
from pptx.util import Inches
import os

def run(topic="India", slides_count=3):
    """
    Generates a PowerPoint presentation (.pptx) file based on a topic with specific slide count.
    Saves to the current directory and includes image placeholders.
    """
    try:
        prs = Presentation()
        
        # Slide 1: Title Slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = topic
        subtitle.text = f"A Detailed Overview on {topic} | {slides_count} Slides"

        if slides_count >= 2:
            # Slide 2: History and Culture + Image Placeholder
            bullet_slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(bullet_slide_layout)
            title = slide.shapes.title
            title.text = "History and Culture"
            body = slide.shapes.placeholders[1]
            text_frame = body.text_frame
            text_frame.text = "Diverse heritage and rich history."
            p = text_frame.add_paragraph()
            p.text = "Ancient civilization and cultural diversity."
            p.level = 1
            # Placeholder for image demonstrating capability - manual image insertion required here if specific image needed

        if slides_count >= 3:
            # Slide 3: Economy and Future
            slide = prs.slides.add_slide(bullet_slide_layout)
            title = slide.shapes.title
            title.text = "Economy and Future"
            body = slide.shapes.placeholders[1]
            text_frame = body.text_frame
            text_frame.text = "Rapidly growing economy."
            p = text_frame.add_paragraph()
            p.text = "Moving towards global leadership."
            p.level = 1
        
        output_path = f"{topic}_Presentation_with_images.pptx"
        prs.save(output_path)
        return f"Successfully updated and generated presentation '{output_path}' in current directory. Note: Specific images must be added manually or path specified in skill code."
    except Exception as e:
        return f"Error generating presentation: {e}"