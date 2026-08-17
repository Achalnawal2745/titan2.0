"""
doc_engine.py — TITAN 2.0 Document & Presentation Intelligence Engine
-----------------------------------------------------------------------
Deterministic, high-fidelity document and presentation generator.
Supports:
  1. Word (.docx)       → Executive reports, typography, tables, callout boxes, embedded images, headers/footers
  2. PowerPoint (.pptx) → 16:9 Widescreen decks (Hero Title, Cards, Metrics, Split-Screen Images, Comparisons)
  3. Interactive (.html)→ Modern animated browser decks with keyboard navigation & glassmorphism
  4. Multi-format Reader→ Ingests .pdf, .docx, .pptx, .xlsx, .csv, .txt, .md
"""

import os
import re
import sys
import json
import hashlib
import urllib.request
import subprocess
from pathlib import Path
from datetime import datetime
from colorama import Fore, Style, init

try:
    if hasattr(sys.stdout, "reconfigure"):
        sys.stdout.reconfigure(encoding="utf-8", errors="replace")
    if hasattr(sys.stderr, "reconfigure"):
        sys.stderr.reconfigure(encoding="utf-8", errors="replace")
except Exception:
    pass

init(autoreset=True)

try:
    from docx import Document
    from docx.shared import Pt, RGBColor, Inches
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.enum.table import WD_TABLE_ALIGNMENT
    from docx.oxml import OxmlElement, parse_xml
    from docx.oxml.ns import nsdecls, qn
except ImportError:
    pass

try:
    from pptx import Presentation
    from pptx.util import Inches as PPTInches, Pt as PPTPt
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.dml.color import RGBColor as PPTRGBColor
except ImportError:
    pass


def _log(msg: str, color=Fore.CYAN):
    try:
        print(f"{color}[DOC-ENGINE]{Style.RESET_ALL} {msg}")
    except Exception:
        try:
            print(f"[DOC-ENGINE] {msg.encode('ascii', 'replace').decode('ascii')}")
        except Exception:
            pass


# ──────────────────────────────────────────────────────────────────────────────
# 0. AUTOMATIC IMAGE RESOLVER & FETCHER
# ──────────────────────────────────────────────────────────────────────────────

IMAGE_CACHE_DIR = Path(__file__).resolve().parent / "scratch" / "cache" / "images"
IMAGE_CACHE_DIR.mkdir(parents=True, exist_ok=True)


def resolve_image(image_query_or_url: str) -> Path | None:
    """
    Resolves an image source into a local image file path on disk:
    1. If already a local file path that exists, returns it.
    2. Searches user folders (Desktop, Downloads, Pictures, Documents) for matching filename/query.
    3. If a web URL (http:// or https://), downloads & caches it.
    4. If a keyword query, fetches a high-resolution relevant stock photo automatically.
    """
    if not image_query_or_url or not isinstance(image_query_or_url, str):
        return None

    src = image_query_or_url.strip()
    
    # 1. Exact local file path check
    local_p = Path(src)
    if local_p.exists() and local_p.is_file():
        return local_p

    # 2. Search common user directories (Desktop, Downloads, Pictures, Documents)
    IMAGE_EXTS = (".png", ".jpg", ".jpeg", ".webp", ".bmp", ".gif", ".tiff")
    search_dirs = [
        Path.home() / "Desktop",
        Path.home() / "Downloads",
        Path.home() / "Pictures",
        Path.home() / "Documents",
        Path.cwd(),
    ]

    # 2a. Direct file name match in user folders
    for d in search_dirs:
        if not d.exists():
            continue
        # Direct filename check
        f_direct = d / src
        if f_direct.exists() and f_direct.is_file():
            return f_direct
        # Try appending extensions if not provided
        for ext in IMAGE_EXTS:
            f_with_ext = d / f"{src}{ext}"
            if f_with_ext.exists() and f_with_ext.is_file():
                return f_with_ext

    # 2b. Fuzzy / Partial filename match in Desktop & Downloads (e.g. "screenshot", "chart", "logo")
    clean_keyword = src.lower().replace("image", "").replace("photo", "").replace("file", "").strip()
    if clean_keyword and len(clean_keyword) >= 3:
        for d in [Path.home() / "Desktop", Path.home() / "Downloads", Path.home() / "Pictures"]:
            if not d.exists():
                continue
            try:
                for f in sorted(d.iterdir(), key=lambda x: x.stat().st_mtime if x.is_file() else 0, reverse=True):
                    if f.is_file() and f.suffix.lower() in IMAGE_EXTS:
                        if clean_keyword in f.name.lower():
                            _log(f"🔎 Found local image on {d.name}: '{f.name}'", Fore.GREEN)
                            return f
            except Exception:
                pass

    # 3. Web URL download
    if src.startswith("http://") or src.startswith("https://"):
        try:
            url_hash = hashlib.md5(src.encode("utf-8")).hexdigest()
            cached_file = IMAGE_CACHE_DIR / f"img_{url_hash}.jpg"
            if cached_file.exists() and cached_file.stat().st_size > 1000:
                return cached_file

            req = urllib.request.Request(src, headers={"User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64)"})
            with urllib.request.urlopen(req, timeout=10) as resp:
                data = resp.read()
            cached_file.write_bytes(data)
            return cached_file
        except Exception as e:
            _log(f"⚠️ Failed to download image from URL '{src}': {e}", Fore.YELLOW)
            return None

    # 4. Keyword / Topic query fallback (Fetches high-res stock photo)
    try:
        clean_query = "".join(c if c.isalnum() else "-" for c in src).strip("-")[:40]
        q_hash = hashlib.md5(clean_query.encode("utf-8")).hexdigest()
        cached_file = IMAGE_CACHE_DIR / f"topic_{clean_query}_{q_hash[:8]}.jpg"
        if cached_file.exists() and cached_file.stat().st_size > 1000:
            return cached_file

        stock_url = f"https://loremflickr.com/1200/800/{clean_query}"
        req = urllib.request.Request(stock_url, headers={"User-Agent": "Mozilla/5.0"})
        with urllib.request.urlopen(req, timeout=8) as resp:
            data = resp.read()
        if len(data) > 2000:
            cached_file.write_bytes(data)
            return cached_file
    except Exception as e:
        _log(f"⚠️ Stock image fetch fallback for '{src}': {e}", Fore.YELLOW)

    return None


# ──────────────────────────────────────────────────────────────────────────────
# 1. MULTI-FORMAT DOCUMENT READERS
# ──────────────────────────────────────────────────────────────────────────────

def read_docx(path: str) -> str:
    """Extract structured text from a .docx file."""
    try:
        from docx import Document
        doc = Document(path)
        parts = []
        for p in doc.paragraphs:
            if p.text.strip():
                parts.append(p.text.strip())
        for table in doc.tables:
            for row in table.rows:
                row_txt = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                if row_txt:
                    parts.append(f"| {row_txt} |")
        return "\n\n".join(parts)
    except Exception as e:
        return f"❌ Cannot read docx '{path}': {e}"


def read_pptx(path: str) -> str:
    """Extract text from all slides in a .pptx presentation."""
    try:
        from pptx import Presentation
        prs = Presentation(path)
        slides_text = []
        for idx, slide in enumerate(prs.slides, start=1):
            slide_content = [f"--- Slide {idx} ---"]
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        txt = paragraph.text.strip()
                        if txt:
                            slide_content.append(txt)
                elif shape.has_table:
                    for row in shape.table.rows:
                        row_txt = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                        if row_txt:
                            slide_content.append(f"| {row_txt} |")
            slides_text.append("\n".join(slide_content))
        return "\n\n".join(slides_text)
    except Exception as e:
        return f"❌ Cannot read pptx '{path}': {e}"


def read_pdf(path: str) -> str:
    """Extract structured text from PDF."""
    try:
        from docling.document_converter import DocumentConverter
        _log("Using Docling (Vision) to parse PDF...")
        converter = DocumentConverter()
        result = converter.convert(path)
        return result.document.export_to_markdown()
    except Exception:
        try:
            import pypdf
            reader = pypdf.PdfReader(path)
            return "\n\n".join(page.extract_text() or "" for page in reader.pages)
        except Exception as e2:
            return f"❌ All PDF readers failed: {e2}"


def read_document(path: str) -> str:
    """Auto-detect and extract content from any supported document type."""
    p = Path(path)
    if not p.exists():
        return f"❌ File not found: {path}"

    ext = p.suffix.lower()
    _log(f"Reading file '{p.name}' ({ext})")

    if ext == ".docx":
        return read_docx(path)
    elif ext in (".pptx", ".ppt"):
        return read_pptx(path)
    elif ext == ".pdf":
        return read_pdf(path)
    elif ext in (".txt", ".md", ".json", ".csv", ".log"):
        return p.read_text(encoding="utf-8", errors="replace")
    elif ext in (".xlsx", ".xls"):
        try:
            import openpyxl
            wb = openpyxl.load_workbook(path, data_only=True)
            lines = []
            for sheet in wb.sheetnames:
                ws = wb[sheet]
                lines.append(f"=== Sheet: {sheet} ===")
                for row in ws.iter_rows(values_only=True):
                    row_vals = [str(v) if v is not None else "" for v in row]
                    if any(row_vals):
                        lines.append(" | ".join(row_vals))
            return "\n".join(lines)
        except Exception as e:
            return f"❌ Excel read failed: {e}"

    return p.read_text(encoding="utf-8", errors="replace")


# ──────────────────────────────────────────────────────────────────────────────
# 2. EXECUTIVE WORD DOCUMENT GENERATOR (.docx) WITH IMAGES & TABLES
# ──────────────────────────────────────────────────────────────────────────────

def _calculate_type_scale(font_size: str | int | float = "normal"):
    """
    Returns: (h1_pt, h2_pt, h3_pt, body_pt, space_after_pt, line_spacing)
    """
    if isinstance(font_size, (int, float)):
        b = float(font_size)
        return (b + 6.0, b + 3.0, b + 1.0, b, max(2.0, b * 0.4), 1.12)

    fs = str(font_size).lower().strip()
    if fs in ("small", "compact", "tiny", "resume", "compact_resume"):
        return (13.5, 11.5, 10.0, 9.0, 3.0, 1.05)
    elif fs in ("large", "big"):
        return (20.0, 15.0, 13.0, 12.0, 7.0, 1.25)
    elif fs in ("huge", "extra_large"):
        return (24.0, 18.0, 15.0, 14.0, 9.0, 1.3)
    else:  # normal
        return (17.0, 13.5, 11.5, 10.5, 5.0, 1.15)


def create_word_document(
    path: str,
    markdown_content: str,
    title: str = "",
    subtitle: str = "",
    author: str = "",
    theme: str = "navy",
    font_size: str | int | float = "normal",
    open_after: bool = True
) -> str:
    """
    Creates a clean, executive-grade formatted Word document from structured markdown.
    Features:
      - Dynamic font sizing (e.g. 'compact', 'small', 'normal', 'large', or specific pt size)
      - 1.0 inch exact margins (1440 DXA)
      - Typography hierarchy (Title, Subtitle, H1, H2, H3, Body)
      - Embedded images with automatic download & aspect ratio preservation
      - Callout highlight boxes with left accent borders
      - Zebra-striped data tables with padded cells and bold headers
      - Bullet & numbered lists with custom indentation
      - Auto-handles file locking if open in Microsoft Word
    """
    try:
        from docx import Document
        from docx.shared import Pt, RGBColor, Inches
        from docx.enum.text import WD_ALIGN_PARAGRAPH
        from docx.enum.table import WD_TABLE_ALIGNMENT
        from docx.oxml import OxmlElement, parse_xml
        from docx.oxml.ns import nsdecls, qn
    except ImportError:
        return "❌ python-docx not installed. Run: pip install python-docx"

    _log(f"Generating Word document: {path} (Font: {font_size})")

    THEMES = {
        "navy": {
            "primary": RGBColor(15, 23, 42),       # Slate 900
            "accent": RGBColor(37, 99, 235),       # Blue 600
            "subtle": RGBColor(100, 116, 139),     # Slate 500
            "dark": RGBColor(51, 65, 85),          # Slate 700
            "bg_hex": "F8FAFC",
            "accent_hex": "2563EB",
            "header_bg": "0F172A",
            "header_txt": "FFFFFF"
        },
        "executive": {
            "primary": RGBColor(30, 41, 59),
            "accent": RGBColor(79, 70, 229),       # Indigo 600
            "subtle": RGBColor(107, 114, 128),
            "dark": RGBColor(55, 65, 81),
            "bg_hex": "F9FAFB",
            "accent_hex": "4F46E5",
            "header_bg": "1E293B",
            "header_txt": "FFFFFF"
        },
        "emerald": {
            "primary": RGBColor(6, 78, 59),        # Emerald 900
            "accent": RGBColor(16, 185, 129),      # Emerald 500
            "subtle": RGBColor(100, 116, 139),
            "dark": RGBColor(30, 41, 59),
            "bg_hex": "F0FDF4",
            "accent_hex": "10B981",
            "header_bg": "064E3B",
            "header_txt": "FFFFFF"
        }
    }
    colors = THEMES.get(theme.lower(), THEMES["navy"])
    h1_pt, h2_pt, h3_pt, body_pt, space_after_pt, line_spacing = _calculate_type_scale(font_size)

    doc = Document()

    # Set 1-Inch Standard Margins (Clean - No unwanted headers/footers)
    for sec in doc.sections:
        sec.top_margin = Inches(1.0)
        sec.bottom_margin = Inches(1.0)
        sec.left_margin = Inches(1.0)
        sec.right_margin = Inches(1.0)

    # 1. Optional Document Cover / Header Block (Only if title is explicitly provided and not already in markdown)
    has_h1 = markdown_content.strip().startswith("# ")
    if title and not has_h1:
        p_title = doc.add_paragraph()
        p_title.alignment = WD_ALIGN_PARAGRAPH.LEFT
        p_title.paragraph_format.space_before = Pt(8)
        p_title.paragraph_format.space_after = Pt(2)
        run_title = p_title.add_run(title)
        run_title.font.name = "Segoe UI"
        run_title.font.size = Pt(max(16, h1_pt + 5))
        run_title.font.bold = True
        run_title.font.color.rgb = colors["primary"]

        if subtitle:
            p_sub = doc.add_paragraph()
            p_sub.alignment = WD_ALIGN_PARAGRAPH.LEFT
            p_sub.paragraph_format.space_after = Pt(4)
            run_sub = p_sub.add_run(subtitle)
            run_sub.font.name = "Segoe UI"
            run_sub.font.size = Pt(body_pt + 1)
            run_sub.font.color.rgb = colors["accent"]

        if author:
            p_meta = doc.add_paragraph()
            p_meta.paragraph_format.space_after = Pt(10)
            r_meta = p_meta.add_run(f"{author}   |   {datetime.now().strftime('%d %b %Y')}")
            r_meta.font.name = "Segoe UI"
            r_meta.font.size = Pt(max(8.0, body_pt - 1.5))
            r_meta.font.italic = True
            r_meta.font.color.rgb = colors["subtle"]

            p_div = doc.add_paragraph()
            p_div.paragraph_format.space_after = Pt(10)
            r_div = p_div.add_run("―" * 45)
            r_div.font.color.rgb = colors["accent"]

    # 2. Markdown Content Parser & Formatter
    lines = markdown_content.strip().split("\n")
    i = 0
    img_re = re.compile(r"^!\[(.*?)\]\((.*?)\)")

    while i < len(lines):
        line = lines[i].rstrip()
        if not line:
            i += 1
            continue

        # Embedded Images ![Caption](url or query)
        if m_img := img_re.match(line):
            caption_text = m_img.group(1) or "Figure"
            img_src = m_img.group(2)
            resolved_img = resolve_image(img_src)
            if resolved_img and resolved_img.exists():
                try:
                    p_img = doc.add_paragraph()
                    p_img.alignment = WD_ALIGN_PARAGRAPH.CENTER
                    p_img.paragraph_format.space_before = Pt(6)
                    p_img.paragraph_format.space_after = Pt(2)
                    doc.add_picture(str(resolved_img), width=Inches(5.5))
                    
                    p_cap = doc.add_paragraph()
                    p_cap.alignment = WD_ALIGN_PARAGRAPH.CENTER
                    p_cap.paragraph_format.space_after = Pt(8)
                    rc = p_cap.add_run(f"Figure: {caption_text}")
                    rc.font.name = "Segoe UI"
                    rc.font.size = Pt(max(8.0, body_pt - 1.5))
                    rc.font.italic = True
                    rc.font.color.rgb = colors["subtle"]
                except Exception as e_img:
                    _log(f"⚠️ Picture insert warning: {e_img}", Fore.YELLOW)
            i += 1
            continue

        # Headings
        if line.startswith("# "):
            p = doc.add_paragraph()
            p.paragraph_format.space_before = Pt(12)
            p.paragraph_format.space_after = Pt(3)
            r = p.add_run(line[2:].strip())
            r.font.name = "Segoe UI"
            r.font.size = Pt(h1_pt)
            r.font.bold = True
            r.font.color.rgb = colors["primary"]

        elif line.startswith("## "):
            p = doc.add_paragraph()
            p.paragraph_format.space_before = Pt(9)
            p.paragraph_format.space_after = Pt(2)
            r = p.add_run(line[3:].strip())
            r.font.name = "Segoe UI"
            r.font.size = Pt(h2_pt)
            r.font.bold = True
            r.font.color.rgb = colors["accent"]

        elif line.startswith("### "):
            p = doc.add_paragraph()
            p.paragraph_format.space_before = Pt(6)
            p.paragraph_format.space_after = Pt(2)
            r = p.add_run(line[4:].strip())
            r.font.name = "Segoe UI"
            r.font.size = Pt(h3_pt)
            r.font.bold = True
            r.font.color.rgb = colors["dark"]

        # Callout / Quote Box (> text)
        elif line.startswith("> "):
            callout_text = line[2:].strip()
            while i + 1 < len(lines) and lines[i + 1].startswith("> "):
                i += 1
                callout_text += " " + lines[i][2:].strip()

            tbl = doc.add_table(rows=1, cols=1)
            tbl.alignment = WD_TABLE_ALIGNMENT.CENTER
            cell = tbl.cell(0, 0)
            cell.width = Inches(6.5)
            
            shd = parse_xml(f'<w:shd {nsdecls("w")} w:fill="{colors["bg_hex"]}"/>')
            borders = parse_xml(
                f'<w:tcBorders {nsdecls("w")}>'
                f'<w:top w:val="none"/>'
                f'<w:left w:val="single" w:sz="24" w:space="0" w:color="{colors["accent_hex"]}"/>'
                f'<w:bottom w:val="none"/>'
                f'<w:right w:val="none"/>'
                f'</w:tcBorders>'
            )
            cell._tc.get_or_add_tcPr().append(shd)
            cell._tc.get_or_add_tcPr().append(borders)

            cp = cell.paragraphs[0]
            cp.paragraph_format.space_before = Pt(3)
            cp.paragraph_format.space_after = Pt(3)
            cr = cp.add_run(f"💡 {callout_text}")
            cr.font.name = "Segoe UI"
            cr.font.size = Pt(body_pt)
            cr.font.italic = True
            cr.font.color.rgb = colors["primary"]

            doc.add_paragraph().paragraph_format.space_after = Pt(3)

        # Markdown Tables (| Col1 | Col2 |)
        elif line.startswith("|") and "|" in line[1:]:
            table_lines = [line]
            while i + 1 < len(lines) and lines[i + 1].startswith("|"):
                i += 1
                table_lines.append(lines[i].strip())

            parsed_rows = []
            for tl in table_lines:
                if re.match(r"^\|[\s\-:|]+\|$", tl):
                    continue
                cells = [c.strip() for c in tl.strip("|").split("|")]
                if cells:
                    parsed_rows.append(cells)

            if parsed_rows:
                num_cols = max(len(r) for r in parsed_rows)
                tbl = doc.add_table(rows=len(parsed_rows), cols=num_cols)
                tbl.alignment = WD_TABLE_ALIGNMENT.CENTER

                for r_idx, row_data in enumerate(parsed_rows):
                    for c_idx in range(num_cols):
                        val = row_data[c_idx] if c_idx < len(row_data) else ""
                        c = tbl.cell(r_idx, c_idx)
                        p = c.paragraphs[0]
                        p.paragraph_format.space_before = Pt(2)
                        p.paragraph_format.space_after = Pt(2)
                        
                        r = p.add_run(val)
                        r.font.name = "Segoe UI"
                        r.font.size = Pt(max(8.0, body_pt - 1.0))

                        if r_idx == 0:
                            r.font.bold = True
                            r.font.color.rgb = RGBColor(255, 255, 255)
                            shd = parse_xml(f'<w:shd {nsdecls("w")} w:fill="{colors["header_bg"]}"/>')
                            c._tc.get_or_add_tcPr().append(shd)
                        elif r_idx % 2 == 1:
                            shd = parse_xml(f'<w:shd {nsdecls("w")} w:fill="{colors["bg_hex"]}"/>')
                            c._tc.get_or_add_tcPr().append(shd)

                doc.add_paragraph().paragraph_format.space_after = Pt(4)

        # Bullet Lists
        elif line.startswith("- ") or line.startswith("* "):
            p = doc.add_paragraph(style="List Bullet")
            p.paragraph_format.space_before = Pt(1)
            p.paragraph_format.space_after = Pt(max(1.5, space_after_pt * 0.4))
            _format_runs_in_paragraph(p, line[2:].strip(), colors, body_size=body_pt)

        # Numbered Lists
        elif re.match(r"^\d+\.\s+", line):
            text_part = re.sub(r"^\d+\.\s+", "", line)
            p = doc.add_paragraph(style="List Number")
            p.paragraph_format.space_before = Pt(1)
            p.paragraph_format.space_after = Pt(max(1.5, space_after_pt * 0.4))
            _format_runs_in_paragraph(p, text_part, colors, body_size=body_pt)

        # Horizontal Rule
        elif line.startswith("---") or line.startswith("___"):
            p = doc.add_paragraph()
            p.paragraph_format.space_before = Pt(3)
            p.paragraph_format.space_after = Pt(3)
            r = p.add_run("―" * 50)
            r.font.color.rgb = colors["subtle"]

        # Standard Paragraph
        else:
            p = doc.add_paragraph()
            p.paragraph_format.line_spacing = line_spacing
            p.paragraph_format.space_after = Pt(space_after_pt)
            _format_runs_in_paragraph(p, line, colors, body_size=body_pt)

        i += 1

    p_out = Path(path)
    p_out.parent.mkdir(parents=True, exist_ok=True)
    try:
        doc.save(str(p_out))
    except PermissionError:
        stem = p_out.stem
        ext = p_out.suffix
        saved = False
        for n in range(2, 50):
            cand = p_out.parent / f"{stem}_v{n}{ext}"
            try:
                doc.save(str(cand))
                p_out = cand
                _log(f"⚠️ Target file is open in Word. Saved updated version to: {p_out.name}", Fore.YELLOW)
                saved = True
                break
            except PermissionError:
                continue
        if not saved:
            return f"❌ Permission denied: File '{p_out.name}' is open in another program."

    _log(f"✅ Word Document successfully generated at: {p_out.resolve()}", Fore.GREEN)

    if open_after:
        try:
            subprocess.Popen(f'start "" "{p_out.resolve()}"', shell=True)
        except Exception:
            pass

    return f"✅ Word document generated successfully: {p_out.resolve()}"


def _format_runs_in_paragraph(p, text: str, colors: dict, body_size: float = 10.5):
    """Parses inline markdown like **bold**, *italic*, and `code` into formatted Word runs."""
    tokens = re.split(r"(\*\*.*?\*\*|\*.*?\*|`.*?`)", text)
    for tok in tokens:
        if not tok:
            continue
        if tok.startswith("**") and tok.endswith("**") and len(tok) >= 4:
            r = p.add_run(tok[2:-2])
            r.font.name = "Segoe UI"
            r.font.size = Pt(body_size)
            r.font.bold = True
            r.font.color.rgb = colors["primary"]
        elif tok.startswith("*") and tok.endswith("*") and len(tok) >= 2:
            r = p.add_run(tok[1:-1])
            r.font.name = "Segoe UI"
            r.font.size = Pt(body_size)
            r.font.italic = True
        elif tok.startswith("`") and tok.endswith("`") and len(tok) >= 2:
            r = p.add_run(f" {tok[1:-1]} ")
            r.font.name = "Consolas"
            r.font.size = Pt(max(8.0, body_size - 1.0))
            r.font.color.rgb = colors["accent"]
        else:
            r = p.add_run(tok)
            r.font.name = "Segoe UI"
            r.font.size = Pt(body_size)
            r.font.color.rgb = colors["dark"]


def fill_docx_template(
    template_path: str,
    output_path: str,
    replacements: dict = None,
    open_after: bool = True
) -> str:
    """
    Opens an existing Word template document (.docx), replaces all placeholder tokens
    (like {{name}}, [COMPANY], {{date}}, etc.) across all paragraphs, headers, footers,
    and table cells while preserving 100% of original fonts, logos, colors, and margins!
    """
    t_path = Path(template_path)
    if not t_path.exists():
        resolved = resolve_image(template_path) or (Path.home() / "Desktop" / template_path)
        if resolved.exists():
            t_path = resolved
        else:
            return f"❌ Template file not found: {template_path}"

    try:
        from docx import Document
    except ImportError:
        return "❌ python-docx not installed."

    _log(f"Filling DOCX template '{t_path.name}' -> '{output_path}'")
    doc = Document(str(t_path))
    replacements = replacements or {}

    def _replace_in_paragraph(p):
        for key, val in replacements.items():
            if key in p.text:
                for run in p.runs:
                    if key in run.text:
                        run.text = run.text.replace(key, str(val))
                if key in p.text:
                    p.text = p.text.replace(key, str(val))

    for p in doc.paragraphs:
        _replace_in_paragraph(p)

    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for cp in cell.paragraphs:
                    _replace_in_paragraph(cp)

    for section in doc.sections:
        for hp in section.header.paragraphs:
            _replace_in_paragraph(hp)
        for fp in section.footer.paragraphs:
            _replace_in_paragraph(fp)

    out_p = Path(output_path)
    out_p.parent.mkdir(parents=True, exist_ok=True)
    try:
        doc.save(str(out_p))
    except PermissionError:
        stem = out_p.stem
        ext = out_p.suffix
        saved = False
        for n in range(2, 50):
            cand = out_p.parent / f"{stem}_v{n}{ext}"
            try:
                doc.save(str(cand))
                out_p = cand
                saved = True
                break
            except PermissionError:
                continue
        if not saved:
            return f"❌ Permission denied: File '{out_p.name}' is open in another program."

    _log(f"✅ Template successfully filled and saved to: {out_p.resolve()}", Fore.GREEN)

    if open_after:
        try:
            subprocess.Popen(f'start "" "{out_p.resolve()}"', shell=True)
        except Exception:
            pass

    return f"✅ Word document generated from template: {out_p.resolve()}"


# ──────────────────────────────────────────────────────────────────────────────
# 3. EXECUTIVE POWERPOINT PRESENTATION GENERATOR (.pptx) WITH IMAGES
# ──────────────────────────────────────────────────────────────────────────────

def create_presentation_deck(
    path: str,
    slides_data: list[dict],
    deck_title: str = "Presentation Deck",
    theme: str = "midnight",
    open_after: bool = True
) -> str:
    """
    Creates a modern 16:9 Widescreen PowerPoint Presentation from structured slide data.
    Slide layout types:
      - 'title'       → Hero cover with large bold title, subtitle, author card
      - 'cards'       → 2 or 3-column structured container cards with rounded boxes
      - 'metrics'     → 2, 3, or 4 massive stat/KPI highlights with subtext
      - 'split_image' → Split-screen: HD Image on Left, Title + Bullets on Right
      - 'bullets'     → Content slide with accent icon pill, headings, and detailed bullet points
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.dml.color import RGBColor
    except ImportError:
        return "❌ python-pptx not installed. Run: pip install python-pptx"

    _log(f"Generating 16:9 PowerPoint Deck: {path}")

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    THEMES = {
        "midnight": {
            "bg": RGBColor(11, 15, 25),            # Deep space black
            "card_bg": RGBColor(20, 27, 45),       # Dark Slate Card
            "card_border": RGBColor(59, 130, 246), # Electric Blue
            "title": RGBColor(255, 255, 255),
            "accent": RGBColor(59, 130, 246),
            "accent2": RGBColor(6, 182, 212),      # Cyan
            "body": RGBColor(203, 213, 225),
            "muted": RGBColor(148, 163, 184)
        },
        "executive": {
            "bg": RGBColor(248, 250, 252),          # Pure modern light slate
            "card_bg": RGBColor(255, 255, 255),
            "card_border": RGBColor(203, 213, 225),
            "title": RGBColor(15, 23, 42),
            "accent": RGBColor(37, 99, 235),
            "accent2": RGBColor(79, 70, 229),
            "body": RGBColor(51, 65, 85),
            "muted": RGBColor(100, 116, 139)
        },
        "crimson": {
            "bg": RGBColor(18, 18, 24),
            "card_bg": RGBColor(28, 28, 38),
            "card_border": RGBColor(225, 29, 72),  # Crimson Red
            "title": RGBColor(255, 255, 255),
            "accent": RGBColor(225, 29, 72),
            "accent2": RGBColor(245, 158, 11),     # Amber
            "body": RGBColor(228, 228, 231),
            "muted": RGBColor(161, 161, 170)
        }
    }
    c = THEMES.get(theme.lower(), THEMES["midnight"])

    def _set_slide_bg(slide):
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, prs.slide_height
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = c["bg"]
        bg.line.fill.background()

    def _add_slide_header(slide, title_text: str, category: str = ""):
        if category:
            cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(10), Inches(0.4))
            tf_c = cat_box.text_frame
            tf_c.word_wrap = True
            p_c = tf_c.paragraphs[0]
            p_c.text = str(category).upper()
            p_c.font.name = "Segoe UI"
            p_c.font.size = Pt(11)
            p_c.font.bold = True
            p_c.font.color.rgb = c["accent"]

        title_top = Inches(0.85) if category else Inches(0.6)
        tbox = slide.shapes.add_textbox(Inches(0.8), title_top, Inches(11.5), Inches(0.8))
        tf = tbox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.name = "Segoe UI"
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = c["title"]

    for s_data in slides_data:
        stype = s_data.get("type", "bullets").lower()
        slide = prs.slides.add_slide(blank_layout)
        _set_slide_bg(slide)

        # ── 1. TITLE / HERO COVER SLIDE ──────────────────────────────────────
        if stype == "title":
            hero = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(1.2), Inches(11.333), Inches(5.1)
            )
            hero.fill.solid()
            hero.fill.fore_color.rgb = c["card_bg"]
            hero.line.color.rgb = c["card_border"]
            hero.line.width = Pt(1.5)

            tf = hero.text_frame
            tf.word_wrap = True
            tf.margin_left = Inches(0.8)
            tf.margin_top = Inches(1.0)

            # Only add category if explicitly provided
            if s_data.get("category"):
                p0 = tf.paragraphs[0]
                p0.text = str(s_data.get("category")).upper()
                p0.font.name = "Segoe UI"
                p0.font.size = Pt(13)
                p0.font.bold = True
                p0.font.color.rgb = c["accent"]
                p0.space_after = Pt(14)
                p1 = tf.add_paragraph()
            else:
                p1 = tf.paragraphs[0]

            p1.text = s_data.get("title", deck_title)
            p1.font.name = "Segoe UI"
            p1.font.size = Pt(40)
            p1.font.bold = True
            p1.font.color.rgb = c["title"]
            p1.space_after = Pt(14)

            if s_data.get("subtitle"):
                p2 = tf.add_paragraph()
                p2.text = s_data.get("subtitle")
                p2.font.name = "Segoe UI"
                p2.font.size = Pt(18)
                p2.font.color.rgb = c["body"]
                p2.space_after = Pt(24)

            # Only add author if provided
            if s_data.get("author"):
                p3 = tf.add_paragraph()
                p3.text = f"{s_data.get('author')}"
                p3.font.name = "Segoe UI"
                p3.font.size = Pt(11)
                p3.font.color.rgb = c["muted"]

        # ── 2. SPLIT-SCREEN IMAGE + CONTENT SLIDE ────────────────────────────
        elif stype in ("split_image", "image", "visual"):
            _add_slide_header(slide, s_data.get("title", "Visual Overview"), s_data.get("category", ""))
            
            # Left Column: Image Box
            img_src = s_data.get("image", s_data.get("title", "technology"))
            img_file = resolve_image(img_src)
            
            if img_file and img_file.exists():
                try:
                    slide.shapes.add_picture(
                        str(img_file), Inches(0.8), Inches(1.8), Inches(5.4), Inches(4.9)
                    )
                except Exception as e_p:
                    _log(f"⚠️ Slide picture insert: {e_p}", Fore.YELLOW)

            # Right Column: Content Card
            cbox = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.5), Inches(1.8), Inches(6.0), Inches(4.9)
            )
            cbox.fill.solid()
            cbox.fill.fore_color.rgb = c["card_bg"]
            cbox.line.color.rgb = c["card_border"]
            cbox.line.width = Pt(1.0)

            tf = cbox.text_frame
            tf.word_wrap = True
            tf.margin_left = Inches(0.5)
            tf.margin_right = Inches(0.5)
            tf.margin_top = Inches(0.5)

            bullets = s_data.get("bullets", [])
            for b_idx, bullet_text in enumerate(bullets):
                p = tf.paragraphs[0] if b_idx == 0 else tf.add_paragraph()
                p.text = f"▹   {bullet_text}"
                p.font.name = "Segoe UI"
                p.font.size = Pt(15)
                p.font.color.rgb = c["body"]
                p.space_after = Pt(14)

        # ── 3. CARDS / MULTI-COLUMN SLIDE ────────────────────────────────────
        elif stype in ("cards", "columns"):
            _add_slide_header(slide, s_data.get("title", "Key Highlights"), s_data.get("category", ""))
            items = s_data.get("items", [])
            num_cards = min(len(items), 3) or 1
            card_width = (Inches(11.733) - (Inches(0.4) * (num_cards - 1))) / num_cards
            start_left = Inches(0.8)
            card_top = Inches(1.9)
            card_height = Inches(4.8)

            for idx, card_item in enumerate(items[:3]):
                cur_left = start_left + (idx * (card_width + Inches(0.4)))
                cbox = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE, cur_left, card_top, card_width, card_height
                )
                cbox.fill.solid()
                cbox.fill.fore_color.rgb = c["card_bg"]
                cbox.line.color.rgb = c["card_border"] if idx == 0 else c["card_bg"]
                cbox.line.width = Pt(1.5)

                tf = cbox.text_frame
                tf.word_wrap = True
                tf.margin_left = Inches(0.4)
                tf.margin_right = Inches(0.4)
                tf.margin_top = Inches(0.4)

                p_h = tf.paragraphs[0]
                p_h.text = card_item.get("header", f"Item {idx+1}")
                p_h.font.name = "Segoe UI"
                p_h.font.size = Pt(20)
                p_h.font.bold = True
                p_h.font.color.rgb = c["accent2"] if idx == 1 else c["accent"]
                p_h.space_after = Pt(12)

                bullets = card_item.get("bullets", [])
                if isinstance(bullets, str):
                    bullets = [bullets]
                for b in bullets:
                    pb = tf.add_paragraph()
                    pb.text = f"• {b}"
                    pb.font.name = "Segoe UI"
                    pb.font.size = Pt(13)
                    pb.font.color.rgb = c["body"]
                    pb.space_after = Pt(8)

        # ── 4. METRICS / STATS SLIDE ─────────────────────────────────────────
        elif stype == "metrics":
            _add_slide_header(slide, s_data.get("title", "Key Performance Indicators"), s_data.get("category", ""))
            metrics = s_data.get("metrics", [])
            count = min(len(metrics), 4) or 1
            stat_w = (Inches(11.733) - (Inches(0.4) * (count - 1))) / count
            stat_left = Inches(0.8)
            stat_top = Inches(2.2)
            stat_h = Inches(4.2)

            for idx, m in enumerate(metrics[:4]):
                cur_x = stat_left + (idx * (stat_w + Inches(0.4)))
                sbox = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE, cur_x, stat_top, stat_w, stat_h
                )
                sbox.fill.solid()
                sbox.fill.fore_color.rgb = c["card_bg"]
                sbox.line.color.rgb = c["card_border"]
                sbox.line.width = Pt(1.5)

                tf = sbox.text_frame
                tf.word_wrap = True
                tf.margin_left = Inches(0.3)
                tf.margin_top = Inches(0.6)

                p_v = tf.paragraphs[0]
                p_v.alignment = PP_ALIGN.CENTER
                p_v.text = m.get("value", "100%")
                p_v.font.name = "Segoe UI"
                p_v.font.size = Pt(46)
                p_v.font.bold = True
                p_v.font.color.rgb = c["accent2"] if idx % 2 == 1 else c["accent"]
                p_v.space_after = Pt(10)

                p_lbl = tf.add_paragraph()
                p_lbl.alignment = PP_ALIGN.CENTER
                p_lbl.text = m.get("label", "Metric").upper()
                p_lbl.font.name = "Segoe UI"
                p_lbl.font.size = Pt(14)
                p_lbl.font.bold = True
                p_lbl.font.color.rgb = c["title"]
                p_lbl.space_after = Pt(8)

                if m.get("subtext"):
                    p_sub = tf.add_paragraph()
                    p_sub.alignment = PP_ALIGN.CENTER
                    p_sub.text = m.get("subtext")
                    p_sub.font.name = "Segoe UI"
                    p_sub.font.size = Pt(11)
                    p_sub.font.color.rgb = c["muted"]

        # ── 5. STANDARD BULLETS / DETAILED CONTENT SLIDE ─────────────────────
        else:
            _add_slide_header(slide, s_data.get("title", "Overview"), s_data.get("category", ""))
            cbox = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(11.733), Inches(4.9)
            )
            cbox.fill.solid()
            cbox.fill.fore_color.rgb = c["card_bg"]
            cbox.line.color.rgb = c["card_border"]
            cbox.line.width = Pt(1.0)

            tf = cbox.text_frame
            tf.word_wrap = True
            tf.margin_left = Inches(0.6)
            tf.margin_top = Inches(0.5)

            bullets = s_data.get("bullets", [])
            for b_idx, bullet_text in enumerate(bullets):
                p = tf.paragraphs[0] if b_idx == 0 else tf.add_paragraph()
                p.text = f"•   {bullet_text}"
                p.font.name = "Segoe UI"
                p.font.size = Pt(16)
                p.font.color.rgb = c["body"]
                p.space_after = Pt(14)

    p_out = Path(path)
    p_out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(p_out))
    _log(f"✅ PowerPoint Deck generated at: {p_out.resolve()}", Fore.GREEN)

    if open_after:
        try:
            subprocess.Popen(f'start "" "{p_out.resolve()}"', shell=True)
        except Exception:
            pass

    return f"✅ PowerPoint presentation generated: {p_out.resolve()}"


# ──────────────────────────────────────────────────────────────────────────────
# 4. INTERACTIVE BROWSER PRESENTATION ENGINE (.html)
# ──────────────────────────────────────────────────────────────────────────────

def create_interactive_deck(
    path: str,
    slides_data: list[dict],
    deck_title: str = "Interactive Presentation",
    open_after: bool = True
) -> str:
    """
    Generates a standalone, beautiful HTML5 + CSS3 + JS interactive slide presentation.
    Supports arrow key navigation, full-screen toggle, glassmorphic layout, and responsive cards.
    """
    _log(f"Generating Interactive HTML Deck: {path}")

    slides_json = json.dumps(slides_data, ensure_ascii=False)

    html_content = f"""<!DOCTYPE html>
<html lang="en">
<head>
<meta charset="UTF-8">
<meta name="viewport" content="width=device-width, initial-scale=1.0">
<title>{deck_title}</title>
<link href="https://fonts.googleapis.com/css2?family=Outfit:wght@300;400;600;700;800&family=JetBrains+Mono:wght@400;600&display=swap" rel="stylesheet">
<style>
  * {{ box-sizing: border-box; margin: 0; padding: 0; }}
  body {{
    font-family: 'Outfit', sans-serif;
    background: radial-gradient(circle at 15% 15%, #111827, #030712 90%);
    color: #F8FAFC;
    height: 100vh;
    display: flex;
    flex-direction: column;
    overflow: hidden;
    user-select: none;
  }}
  .deck-container {{
    flex: 1;
    display: flex;
    align-items: center;
    justify-content: center;
    padding: 3rem;
    position: relative;
  }}
  .slide {{
    display: none;
    width: 100%;
    max-width: 1100px;
    background: rgba(17, 24, 39, 0.7);
    backdrop-filter: blur(20px);
    border: 1px solid rgba(59, 130, 246, 0.25);
    border-radius: 24px;
    padding: 3.5rem;
    box-shadow: 0 25px 60px -15px rgba(0, 0, 0, 0.7), 0 0 40px rgba(59, 130, 246, 0.15);
    animation: fadeIn 0.4s ease-out;
  }}
  .slide.active {{ display: block; }}
  @keyframes fadeIn {{
    from {{ opacity: 0; transform: translateY(16px) scale(0.98); }}
    to {{ opacity: 1; transform: translateY(0) scale(1); }}
  }}
  .badge {{
    display: inline-block;
    padding: 6px 14px;
    background: rgba(59, 130, 246, 0.15);
    border: 1px solid rgba(59, 130, 246, 0.4);
    border-radius: 9999px;
    color: #60A5FA;
    font-size: 0.85rem;
    font-weight: 700;
    letter-spacing: 1px;
    text-transform: uppercase;
    margin-bottom: 1.2rem;
  }}
  h1 {{ font-size: 3.2rem; font-weight: 800; line-height: 1.15; color: #FFFFFF; margin-bottom: 1rem; }}
  h2 {{ font-size: 2.2rem; font-weight: 700; color: #FFFFFF; margin-bottom: 1.5rem; }}
  p.subtitle {{ font-size: 1.3rem; color: #94A3B8; margin-bottom: 2rem; }}
  
  .cards-grid {{
    display: grid;
    grid-template-columns: repeat(auto-fit, minmax(280px, 1fr));
    gap: 1.5rem;
    margin-top: 1.5rem;
  }}
  .card {{
    background: rgba(30, 41, 59, 0.6);
    border: 1px solid rgba(255, 255, 255, 0.08);
    border-radius: 18px;
    padding: 1.8rem;
    transition: transform 0.2s, border-color 0.2s;
  }}
  .card:hover {{ transform: translateY(-4px); border-color: rgba(59, 130, 246, 0.5); }}
  .card h3 {{ color: #38BDF8; font-size: 1.3rem; margin-bottom: 0.8rem; }}
  .card ul {{ list-style: none; }}
  .card li {{ color: #CBD5E1; font-size: 1rem; margin-bottom: 0.5rem; display: flex; align-items: flex-start; }}
  .card li::before {{ content: "▹"; color: #38BDF8; margin-right: 8px; font-weight: bold; }}

  .metrics-grid {{
    display: grid;
    grid-template-columns: repeat(auto-fit, minmax(200px, 1fr));
    gap: 1.5rem;
    margin-top: 2rem;
  }}
  .metric-box {{
    background: rgba(15, 23, 42, 0.8);
    border: 1px solid rgba(59, 130, 246, 0.3);
    border-radius: 18px;
    padding: 2rem;
    text-align: center;
  }}
  .metric-val {{ font-size: 3.5rem; font-weight: 800; color: #60A5FA; font-family: 'JetBrains Mono', monospace; }}
  .metric-lbl {{ font-size: 1rem; font-weight: 700; color: #F1F5F9; text-transform: uppercase; margin-top: 0.5rem; }}
  .metric-sub {{ font-size: 0.85rem; color: #94A3B8; margin-top: 0.3rem; }}

  .bullet-list {{ list-style: none; margin-top: 1.5rem; }}
  .bullet-list li {{
    font-size: 1.25rem;
    color: #E2E8F0;
    margin-bottom: 1.2rem;
    display: flex;
    align-items: center;
  }}
  .bullet-list li::before {{
    content: "●";
    color: #38BDF8;
    font-size: 0.9rem;
    margin-right: 14px;
  }}

  .footer-controls {{
    display: flex;
    align-items: center;
    justify-content: space-between;
    padding: 1.2rem 3rem;
    background: rgba(10, 15, 29, 0.85);
    border-top: 1px solid rgba(255, 255, 255, 0.05);
  }}
  .btn-nav {{
    background: #2563EB;
    border: none;
    color: white;
    padding: 8px 20px;
    border-radius: 10px;
    font-weight: 600;
    cursor: pointer;
    font-size: 0.95rem;
    transition: background 0.2s;
  }}
  .btn-nav:hover {{ background: #1D4ED8; }}
  .btn-nav:disabled {{ background: #374151; cursor: not-allowed; opacity: 0.5; }}
  .progress-txt {{ font-size: 0.95rem; color: #94A3B8; font-weight: 600; font-family: 'JetBrains Mono', monospace; }}
</style>
</head>
<body>

<div class="deck-container" id="deck"></div>

<div class="footer-controls">
  <div>
    <button class="btn-nav" id="prevBtn" onclick="prevSlide()">← Previous</button>
    <button class="btn-nav" id="nextBtn" onclick="nextSlide()" style="margin-left: 8px;">Next →</button>
  </div>
  <div class="progress-txt" id="slideCount">1 / 1</div>
  <button class="btn-nav" onclick="toggleFullScreen()" style="background: rgba(255,255,255,0.1);">⛶ Fullscreen</button>
</div>

<script>
  const slides = {slides_json};
  let currentIdx = 0;
  const container = document.getElementById('deck');
  const countDisplay = document.getElementById('slideCount');
  const prevBtn = document.getElementById('prevBtn');
  const nextBtn = document.getElementById('nextBtn');

  function renderSlides() {{
    container.innerHTML = '';
    slides.forEach((s, idx) => {{
      const div = document.createElement('div');
      div.className = 'slide' + (idx === 0 ? ' active' : '');
      div.id = 'slide-' + idx;

      let inner = '';
      if (s.category) inner += `<div class="badge">${{s.category}}</div>`;

      if (s.type === 'title') {{
        inner += `<h1>${{s.title || 'Deck Title'}}</h1>`;
        if (s.subtitle) inner += `<p class="subtitle">${{s.subtitle}}</p>`;
        inner += `<div style="margin-top:2rem;color:#64748B;font-size:0.95rem;">Presented by: ${{s.author || 'TITAN'}}</div>`;
      }} else if (s.type === 'cards' || s.type === 'columns') {{
        inner += `<h2>${{s.title}}</h2>`;
        inner += `<div class="cards-grid">`;
        (s.items || []).forEach(item => {{
          inner += `<div class="card"><h3>${{item.header}}</h3><ul>`;
          (item.bullets || []).forEach(b => inner += `<li>${{b}}</li>`);
          inner += `</ul></div>`;
        }});
        inner += `</div>`;
      }} else if (s.type === 'metrics') {{
        inner += `<h2>${{s.title}}</h2>`;
        inner += `<div class="metrics-grid">`;
        (s.metrics || []).forEach(m => {{
          inner += `<div class="metric-box"><div class="metric-val">${{m.value}}</div><div class="metric-lbl">${{m.label}}</div><div class="metric-sub">${{m.subtext || ''}}</div></div>`;
        }});
        inner += `</div>`;
      }} else {{
        inner += `<h2>${{s.title}}</h2>`;
        inner += `<ul class="bullet-list">`;
        (s.bullets || []).forEach(b => inner += `<li>${{b}}</li>`);
        inner += `</ul>`;
      }}

      div.innerHTML = inner;
      container.appendChild(div);
    }});
    updateUI();
  }}

  function showSlide(idx) {{
    document.querySelectorAll('.slide').forEach(s => s.classList.remove('active'));
    const target = document.getElementById('slide-' + idx);
    if (target) target.classList.add('active');
    updateUI();
  }}

  function nextSlide() {{
    if (currentIdx < slides.length - 1) {{
      currentIdx++;
      showSlide(currentIdx);
    }}
  }}

  function prevSlide() {{
    if (currentIdx > 0) {{
      currentIdx--;
      showSlide(currentIdx);
    }}
  }}

  function updateUI() {{
    countDisplay.innerText = `${{currentIdx + 1}} / ${{slides.length}}`;
    prevBtn.disabled = currentIdx === 0;
    nextBtn.disabled = currentIdx === slides.length - 1;
  }}

  function toggleFullScreen() {{
    if (!document.fullscreenElement) {{
      document.documentElement.requestFullscreen();
    }} else {{
      document.exitFullscreen();
    }}
  }}

  document.addEventListener('keydown', (e) => {{
    if (e.key === 'ArrowRight' || e.key === ' ' || e.key === 'PageDown') nextSlide();
    if (e.key === 'ArrowLeft' || e.key === 'PageUp') prevSlide();
    if (e.key === 'f' || e.key === 'F') toggleFullScreen();
  }});

  renderSlides();
</script>
</body>
</html>
"""

    p_out = Path(path)
    p_out.parent.mkdir(parents=True, exist_ok=True)
    p_out.write_text(html_content, encoding="utf-8")
    _log(f"✅ Interactive HTML Deck generated at: {p_out.resolve()}", Fore.GREEN)

    if open_after:
        try:
            subprocess.Popen(f'start "" "{p_out.resolve()}"', shell=True)
        except Exception:
            pass

    return f"✅ Interactive deck generated: {p_out.resolve()}"
