import os
import re
import shutil
import platform
from pathlib import Path
from datetime import datetime

try:
    import send2trash
    _SEND2TRASH = True
except ImportError:
    _SEND2TRASH = False

_OS = platform.system()  # "Windows" | "Darwin" | "Linux"

# On Windows, Downloads/Desktop/Documents/Pictures/Music/Videos are often
# redirected — moved to another drive, or moved into OneDrive by "Known
# Folder Move". Blindly assuming home()/"Downloads" silently looks in the
# WRONG place whenever that's true, which is exactly why "find X in
# Downloads" was failing even though the file was really there.
# We read the actual location from the registry, and only fall back to
# home()/name if that lookup fails (non-Windows, or nothing redirected).
_SHELL_FOLDER_GUIDS = {
    "Desktop":   "{B4BFCC3A-DB2C-424C-B029-7FE99A87C641}",
    "Downloads": "{374DE290-123F-4565-9164-39C4925E467B}",
    "Documents": "{FDD39AD0-238F-46AF-ADB4-6C85480369C7}",
    "Pictures":  "{33E28130-4E1E-4676-835A-98395C3BC3BB}",
    "Music":     "{4BD8D571-6D19-48D3-BE97-422220080E43}",
    "Videos":    "{18989B1D-99B5-455B-841C-AB7C74E4DDFC}",
}
_shell_folder_cache: dict[str, Path] = {}

def _windows_shell_folder(name: str) -> "Path | None":
    """Read the ACTUAL folder path from the registry (handles redirection
    to another drive or OneDrive) instead of assuming it lives under home()."""
    if _OS != "Windows":
        return None
    if name in _shell_folder_cache:
        return _shell_folder_cache[name]
    guid = _SHELL_FOLDER_GUIDS.get(name)
    try:
        import winreg
        for key_name in ("User Shell Folders", "Shell Folders"):
            try:
                with winreg.OpenKey(
                    winreg.HKEY_CURRENT_USER,
                    rf"Software\Microsoft\Windows\CurrentVersion\Explorer\{key_name}",
                ) as key:
                    val = None
                    if guid:
                        try:
                            val, _ = winreg.QueryValueEx(key, guid)
                        except FileNotFoundError:
                            pass
                    if val is None:
                        try:
                            val, _ = winreg.QueryValueEx(key, name)
                        except FileNotFoundError:
                            val = None
                    if val:
                        p = Path(os.path.expandvars(val))
                        if p.exists():
                            _shell_folder_cache[name] = p
                            return p
            except FileNotFoundError:
                continue
    except Exception:
        pass
    return None

_SAFE_ROOTS: list[Path] = [
    Path.home(),
]

def _is_safe_path(target: Path) -> bool:
    """Check if the target path is safe for operations (blocks critical system dirs)."""
    try:
        resolved = target.resolve()
        if _OS == "Windows":
            blocked = {Path("C:/Windows"), Path("C:/Program Files"), Path("C:/Program Files (x86)")}
            if any(resolved == b or resolved.is_relative_to(b) for b in blocked if b.exists()):
                return False
        elif _OS in ("Linux", "Darwin"):
            blocked = {Path("/bin"), Path("/sbin"), Path("/usr/bin"), Path("/usr/sbin"), Path("/etc")}
            if any(resolved == b or resolved.is_relative_to(b) for b in blocked if b.exists()):
                return False
        return True
    except Exception:
        return False

def _get_desktop() -> Path:
    if _OS == "Windows":
        p = _windows_shell_folder("Desktop")
        if p:
            return p
    if _OS == "Linux":
        xdg = os.environ.get("XDG_DESKTOP_DIR", "")
        if xdg and Path(xdg).exists():
            return Path(xdg)
    return Path.home() / "Desktop"

def _get_downloads() -> Path:
    if _OS == "Windows":
        p = _windows_shell_folder("Downloads")
        if p:
            return p
    if _OS == "Linux":
        xdg = os.environ.get("XDG_DOWNLOAD_DIR", "")
        if xdg and Path(xdg).exists():
            return Path(xdg)
    return Path.home() / "Downloads"

def _get_documents() -> Path:
    if _OS == "Windows":
        p = _windows_shell_folder("Documents")
        if p:
            return p
    if _OS == "Linux":
        xdg = os.environ.get("XDG_DOCUMENTS_DIR", "")
        if xdg and Path(xdg).exists():
            return Path(xdg)
    return Path.home() / "Documents"

def _get_pictures() -> Path:
    if _OS == "Windows":
        p = _windows_shell_folder("Pictures")
        if p:
            return p
    if _OS == "Linux":
        xdg = os.environ.get("XDG_PICTURES_DIR", "")
        if xdg and Path(xdg).exists():
            return Path(xdg)
    return Path.home() / "Pictures"

def _get_music() -> Path:
    if _OS == "Windows":
        p = _windows_shell_folder("Music")
        if p:
            return p
    if _OS == "Linux":
        xdg = os.environ.get("XDG_MUSIC_DIR", "")
        if xdg and Path(xdg).exists():
            return Path(xdg)
    return Path.home() / "Music"

def _get_videos() -> Path:
    if _OS == "Windows":
        p = _windows_shell_folder("Videos")
        if p:
            return p
    if _OS == "Linux":
        xdg = os.environ.get("XDG_VIDEOS_DIR", "")
        if xdg and Path(xdg).exists():
            return Path(xdg)
    return Path.home() / "Videos"


def _resolve_path(raw: str) -> Path:
    if not raw or not raw.strip():
        return _get_desktop()

    raw_clean = raw.strip().strip("'\"").replace("\\", "/")
    lower = raw_clean.lower().strip()

    # Handle drive letters and aliases (e.g. "e:", "e:/", "e drive", "drive e")
    for d in "abcdefghijklmnopqrstuvwxyz":
        if lower in (f"{d}:", f"{d}:/", f"{d} drive", f"drive {d}", f"drive {d}:"):
            return Path(f"{d.upper()}:/")

    shortcuts: dict[str, Path] = {
        "desktop":   _get_desktop(),
        "downloads": _get_downloads(),
        "documents": _get_documents(),
        "pictures":  _get_pictures(),
        "music":     _get_music(),
        "videos":    _get_videos(),
        "home":      Path.home(),
    }
    if lower in shortcuts:
        return shortcuts[lower]

    for key, path_obj in shortcuts.items():
        if lower.startswith(key + "/"):
            remainder = raw_clean[len(key) + 1:]
            return path_obj / remainder

    return Path(raw_clean).expanduser()

def _format_size(b: int) -> str:
    for unit in ["B", "KB", "MB", "GB", "TB"]:
        if b < 1024:
            return f"{b:.1f} {unit}"
        b /= 1024
    return f"{b:.1f} TB"

def list_drives() -> str:
    if _OS == "Windows":
        drives = []
        for d in "ABCDEFGHIJKLMNOPQRSTUVWXYZ":
            p = Path(f"{d}:/")
            if p.exists():
                try:
                    usage = shutil.disk_usage(p)
                    free = _format_size(usage.free)
                    total = _format_size(usage.total)
                    drives.append(f"💾 Drive {d}:/ (Free: {free} / Total: {total})")
                except Exception:
                    drives.append(f"💾 Drive {d}:/")
        return "Connected Disks & Drives:\n" + "\n".join(drives) if drives else "No drives detected."
    return "Root filesystem: /"

def _safe_trash(target: Path) -> str:
    if not _SEND2TRASH:
        return (
            "send2trash is not installed. "
            "Run: pip install send2trash — "
            "Permanent deletion is disabled for safety."
        )
    send2trash.send2trash(str(target))
    return f"Moved to Trash: {target.name}"


def list_files(path: str = "desktop", show_hidden: bool = False) -> str:
    try:
        raw_lower = (path or "").strip().lower()
        if raw_lower in ("drives", "all drives", "my computer", "this pc", "disks", "disk"):
            return list_drives()

        target = _resolve_path(path)
        if not _is_safe_path(target):
            return f"Access denied: {target}"
        if not target.exists():
            return f"Path not found: {target}"
        if not target.is_dir():
            return f"Not a directory: {target}"

        items = []
        for item in sorted(target.iterdir()):
            if not show_hidden and item.name.startswith("."):
                continue
            if item.is_dir():
                items.append(f"📁 {item.name}/")
            else:
                size = _format_size(item.stat().st_size)
                items.append(f"📄 {item.name} ({size})")

        display_name = target.name or str(target)
        if not items:
            return f"Directory is empty: {display_name}"

        return f"Contents of {display_name} ({len(items)} items):\n" + "\n".join(items)

    except PermissionError:
        return f"Permission denied: {path}"
    except Exception as e:
        return f"Error listing files: {e}"


def create_file(path: str, name: str = "", content: str = "") -> str:
    try:
        base   = _resolve_path(path)
        target = (base / name) if name else base
        if not _is_safe_path(target):
            return f"Access denied: {target}"
        target.parent.mkdir(parents=True, exist_ok=True)
        target.write_text(content, encoding="utf-8")
        return f"File created: {target.name}"
    except Exception as e:
        return f"Could not create file: {e}"


def create_folder(path: str, name: str = "") -> str:
    try:
        base   = _resolve_path(path)
        target = (base / name) if name else base
        if not _is_safe_path(target):
            return f"Access denied: {target}"
        target.mkdir(parents=True, exist_ok=True)
        return f"Folder created: {target.name}"
    except Exception as e:
        return f"Could not create folder: {e}"


def delete_file(path: str, name: str = "") -> str:
    try:
        base   = _resolve_path(path)
        target = (base / name) if name else base
        if not _is_safe_path(target):
            return f"Access denied: {target}"
        if not target.exists():
            return f"Not found: {target.name}"

        # Güvenli dizin kontrolü — kritik kullanıcı klasörlerini koru
        protected = {
            _get_desktop(), _get_downloads(), _get_documents(),
            _get_pictures(), _get_music(), _get_videos(), Path.home()
        }
        if target.resolve() in {p.resolve() for p in protected}:
            return f"Protected directory, cannot delete: {target.name}"

        return _safe_trash(target)

    except PermissionError:
        return f"Permission denied: {path}"
    except Exception as e:
        return f"Could not delete: {e}"


def move_file(path: str, name: str = "", destination: str = "") -> str:
    try:
        base   = _resolve_path(path)
        src    = (base / name) if name else base
        dst    = _resolve_path(destination) if destination else None

        if not src.exists():
            return f"Source not found: {src.name}"
        if dst is None:
            return "No destination specified."
        if not _is_safe_path(src):
            return f"Access denied (source): {src}"
        if not _is_safe_path(dst):
            return f"Access denied (destination): {dst}"

        if dst.is_dir():
            dst = dst / src.name

        dst.parent.mkdir(parents=True, exist_ok=True)
        shutil.move(str(src), str(dst))
        return f"Moved: {src.name} → {dst.parent.name}/"

    except Exception as e:
        return f"Could not move: {e}"


def copy_file(path: str, name: str = "", destination: str = "") -> str:
    try:
        base = _resolve_path(path)
        src  = (base / name) if name else base
        dst  = _resolve_path(destination) if destination else None

        if not src.exists():
            return f"Source not found: {src.name}"
        if dst is None:
            return "No destination specified."
        if not _is_safe_path(src):
            return f"Access denied (source): {src}"
        if not _is_safe_path(dst):
            return f"Access denied (destination): {dst}"

        if dst.is_dir():
            dst = dst / src.name

        dst.parent.mkdir(parents=True, exist_ok=True)

        if src.is_dir():
            shutil.copytree(str(src), str(dst))
        else:
            shutil.copy2(str(src), str(dst))

        return f"Copied: {src.name} → {dst.parent.name}/"

    except Exception as e:
        return f"Could not copy: {e}"


def rename_file(path: str, name: str = "", new_name: str = "") -> str:
    try:
        base     = _resolve_path(path)
        target   = (base / name) if name else base
        if not _is_safe_path(target):
            return f"Access denied: {target}"
        if not target.exists():
            return f"Not found: {target.name}"
        if not new_name:
            return "No new name provided."

        new_path = target.parent / new_name
        if new_path.exists():
            return f"A file named '{new_name}' already exists here."

        target.rename(new_path)
        return f"Renamed: {target.name} → {new_name}"

    except Exception as e:
        return f"Could not rename: {e}"


def read_file(path: str, name: str = "", max_chars: int = 4000) -> str:
    try:
        base   = _resolve_path(path)
        target = (base / name) if name else base
        if not _is_safe_path(target):
            return f"Access denied: {target}"
        if not target.exists():
            return f"File not found: {target.name}"
        if not target.is_file():
            return f"Not a file: {target.name}"

        ext = target.suffix.lower()
        if ext == ".pdf":
            try:
                import pypdf
                reader = pypdf.PdfReader(str(target))
                pages_text = [page.extract_text() or "" for page in reader.pages]
                content = "\n\n".join(t.strip() for t in pages_text if t.strip())
                if not content:
                    import pdfplumber
                    with pdfplumber.open(str(target)) as pdf:
                        content = "\n\n".join((p.extract_text() or "").strip() for p in pdf.pages if (p.extract_text() or "").strip())
            except Exception:
                try:
                    import pdfplumber
                    with pdfplumber.open(str(target)) as pdf:
                        content = "\n\n".join((p.extract_text() or "").strip() for p in pdf.pages if (p.extract_text() or "").strip())
                except Exception as e:
                    content = f"[PDF text extraction failed: {e}]"

        elif ext in (".docx", ".doc"):
            try:
                from docx import Document
                doc = Document(str(target))
                content = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
            except Exception:
                # OLE2 .doc or corrupted docx text extraction
                try:
                    raw_bytes = target.read_bytes()
                    words = [s.decode("utf-8", "ignore") for s in re.findall(rb"[\x20-\x7e\n\r\t]{4,}", raw_bytes)]
                    content = " ".join(words)
                except Exception as e:
                    content = f"[Document text extraction failed: {e}]"

        elif ext in (".pptx", ".ppt"):
            try:
                from pptx import Presentation
                prs = Presentation(str(target))
                slides = []
                for i, slide in enumerate(prs.slides, 1):
                    txts = [shape.text.strip() for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip()]
                    if txts:
                        slides.append(f"--- Slide {i} ---\n" + "\n".join(txts))
                content = "\n\n".join(slides)
            except Exception as e:
                content = f"[Presentation text extraction failed: {e}]"

        else:
            content = target.read_text(encoding="utf-8", errors="ignore")

        if not content.strip():
            return f"[File '{target.name}' is empty or contains non-extractable media]"

        if len(content) > max_chars:
            content = content[:max_chars] + f"\n\n[Truncated — {len(content)} total chars]"
        return content

    except Exception as e:
        return f"Could not read file: {e}"


def write_file(path: str, name: str = "", content: str = "",
               append: bool = False) -> str:
    try:
        base   = _resolve_path(path)
        target = (base / name) if name else base
        if not _is_safe_path(target):
            return f"Access denied: {target}"
        target.parent.mkdir(parents=True, exist_ok=True)

        # ── Word (.docx) safe writer fix ───────────────────────────────────────
        mode = "a" if append else "w"
        with open(target, mode, encoding="utf-8") as f:
            f.write(content)
        action = "Appended to" if append else "Written to"
        return f"{action}: {target.name}"
    except Exception as e:
        return f"Could not write file: {e}"


def find_files(name: str = "", extension: str = "",
               path: str = "home", max_results: int = 20) -> str:
    try:
        search_path = _resolve_path(path)
        if not _is_safe_path(search_path):
            return f"Access denied: {search_path}"
        if not search_path.exists():
            return f"Search path not found: {path}"

        results    = []
        dir_count  = 0
        max_dirs   = 500  # performans + güvenlik limiti

        for item in search_path.rglob("*"):
            if item.is_dir():
                dir_count += 1
                if dir_count > max_dirs:
                    break
                continue
            if not item.is_file():
                continue
            if extension and item.suffix.lower() != extension.lower():
                continue
            if name and name.lower() not in item.name.lower():
                continue
            size = _format_size(item.stat().st_size)
            results.append(f"📄 {item.name} ({size}) — {item.parent}")
            if len(results) >= max_results:
                break

        if not results:
            query = name or extension or "files"
            return f"No {query} found in {search_path.name}/"

        return f"Found {len(results)} file(s):\n" + "\n".join(results)

    except Exception as e:
        return f"Search error: {e}"


def get_largest_files(path: str = "downloads", count: int = 10) -> str:
    count = min(count, 50)  # maksimum 50
    try:
        search_path = _resolve_path(path)
        if not _is_safe_path(search_path):
            return f"Access denied: {search_path}"
        if not search_path.exists():
            return f"Path not found: {path}"

        files = []
        for item in search_path.rglob("*"):
            if item.is_file():
                try:
                    files.append((item.stat().st_size, item))
                except Exception:
                    continue

        files.sort(reverse=True)
        top = files[:count]

        if not top:
            return "No files found."

        lines = [f"Top {len(top)} largest files in {search_path.name}/:"]
        for size, f in top:
            lines.append(f"  {_format_size(size):>10}  {f.name}  ({f.parent})")

        return "\n".join(lines)

    except Exception as e:
        return f"Error: {e}"


def get_disk_usage(path: str = "home") -> str:
    try:
        target = _resolve_path(path)
        usage  = shutil.disk_usage(target)
        pct    = usage.used / usage.total * 100
        return (
            f"Disk usage ({target}):\n"
            f"  Total : {_format_size(usage.total)}\n"
            f"  Used  : {_format_size(usage.used)} ({pct:.1f}%)\n"
            f"  Free  : {_format_size(usage.free)}"
        )
    except Exception as e:
        return f"Could not get disk usage: {e}"


def organize_desktop() -> str:
    type_map = {
        "Images":    {".jpg", ".jpeg", ".png", ".gif", ".bmp", ".webp", ".svg", ".ico", ".heic"},
        "Documents": {".pdf", ".doc", ".docx", ".txt", ".xls", ".xlsx",
                      ".ppt", ".pptx", ".csv", ".odt", ".ods", ".odp"},
        "Videos":    {".mp4", ".avi", ".mkv", ".mov", ".wmv", ".flv", ".webm", ".m4v"},
        "Music":     {".mp3", ".wav", ".flac", ".aac", ".ogg", ".wma", ".m4a"},
        "Archives":  {".zip", ".rar", ".7z", ".tar", ".gz", ".bz2", ".xz"},
        "Code":      {".py", ".js", ".ts", ".html", ".css", ".json", ".xml",
                      ".cpp", ".java", ".cs", ".go", ".rs", ".sh"},
    }

    desktop = _get_desktop()
    moved, skipped = [], []

    try:
        for item in desktop.iterdir():
            # Klasörlere, gizli dosyalara ve organize klasörlerine dokunma
            if item.is_dir() or item.name.startswith("."):
                continue
            if item.name in {k for k in type_map}:
                continue

            ext        = item.suffix.lower()
            target_dir = desktop / "Others"
            for folder, exts in type_map.items():
                if ext in exts:
                    target_dir = desktop / folder
                    break

            target_dir.mkdir(exist_ok=True)
            new_path = target_dir / item.name

            if new_path.exists():
                skipped.append(item.name)
                continue

            shutil.move(str(item), str(new_path))
            moved.append(f"{item.name} → {target_dir.name}/")

        result = f"Desktop organized: {len(moved)} files moved."
        if moved:
            preview = moved[:8]
            result += "\n" + "\n".join(preview)
            if len(moved) > 8:
                result += f"\n... and {len(moved) - 8} more."
        if skipped:
            result += f"\n{len(skipped)} file(s) skipped (name conflict)."
        return result

    except Exception as e:
        return f"Could not organize desktop: {e}"


def get_file_info(path: str, name: str = "") -> str:
    try:
        base   = _resolve_path(path)
        target = (base / name) if name else base
        if not _is_safe_path(target):
            return f"Access denied: {target}"
        if not target.exists():
            return f"Not found: {target.name}"

        stat = target.stat()
        info = {
            "Name":      target.name,
            "Type":      "Folder" if target.is_dir() else "File",
            "Size":      _format_size(stat.st_size),
            "Location":  str(target.parent),
            "Created":   datetime.fromtimestamp(stat.st_ctime).strftime("%Y-%m-%d %H:%M"),
            "Modified":  datetime.fromtimestamp(stat.st_mtime).strftime("%Y-%m-%d %H:%M"),
            "Extension": target.suffix or "—",
        }
        return "\n".join(f"  {k}: {v}" for k, v in info.items())

    except Exception as e:
        return f"Could not get file info: {e}"

def file_controller(
    parameters: dict = None,
    response=None,
    player=None,
    session_memory=None,
) -> str:
    params = parameters or {}
    action = params.get("action", "").lower().strip()
    path   = (
        params.get("file_path")
        or params.get("filepath")
        or params.get("path")
        or params.get("target")
        or params.get("folder")
        or params.get("directory")
        or "desktop"
    )
    name   = (
        params.get("name")
        or params.get("filename")
        or params.get("file_name")
        or ""
    )
    dest   = (
        params.get("destination")
        or params.get("dest")
        or params.get("to")
        or ""
    )
    n_name = (
        params.get("new_name")
        or params.get("newname")
        or params.get("target_name")
        or ""
    )

    if player:
        player.write_log(f"[file] {action} {name or path}")

    try:
        if action in ("list", "list_files", "ls", "dir"):
            return list_files(path)

        elif action in ("drives", "list_drives"):
            return list_drives()

        elif action in ("create_file", "touch"):
            return create_file(path, name=name, content=params.get("content", ""))

        elif action in ("create_folder", "mkdir"):
            return create_folder(path, name=name)

        elif action in ("delete", "remove", "rm"):
            return delete_file(path, name=name)

        elif action in ("move", "mv"):
            return move_file(path, name=name, destination=dest)

        elif action in ("copy", "cp"):
            return copy_file(path, name=name, destination=dest)

        elif action in ("rename", "ren"):
            return rename_file(path, name=name, new_name=n_name)

        elif action in ("read", "read_file", "view", "cat"):
            return read_file(path, name=name)

        elif action in ("write", "write_file"):
            return write_file(
                path, name=name,
                content=params.get("content", ""),
                append=params.get("append", False)
            )

        elif action in ("find", "search"):
            return find_files(
                name=name or params.get("name", ""),
                extension=params.get("extension", ""),
                path=path,
                max_results=min(int(params.get("max_results", 20)), 50),
            )

        elif action == "largest":
            return get_largest_files(
                path=path,
                count=int(params.get("count", 10)),
            )

        elif action == "disk_usage":
            return get_disk_usage(path)

        elif action == "organize_desktop":
            return organize_desktop()

        elif action in ("info", "stat", "status"):
            return get_file_info(path, name=name)

        else:
            return f"Unknown action: '{action}'"

    except Exception as e:
        return f"File controller error ({action}): {e}"